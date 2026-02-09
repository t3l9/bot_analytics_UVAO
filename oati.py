import os
import re
import time
import asyncio
from datetime import datetime, timedelta
from functools import reduce

import pandas as pd
from dotenv import load_dotenv

# Selenium
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.chrome.service import Service as ChromeService
from webdriver_manager.chrome import ChromeDriverManager

# Telegram
from telegram import (
    Update,
    ReplyKeyboardMarkup,
    InputFile,
    InlineKeyboardButton,
    InlineKeyboardMarkup,
)
from telegram.ext import (
    Application,
    ApplicationBuilder,
    CommandHandler,
    MessageHandler,
    filters,
    ContextTypes,
    ConversationHandler,
    CallbackQueryHandler,
    CallbackContext,
)

# Excel / Windows
from openpyxl import load_workbook
from openpyxl.styles import Border, Side, Alignment, Font, PatternFill
from openpyxl.formatting.rule import CellIsRule
import win32com.client
import pythoncom

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import comtypes.client

try:
    import win32com.client
    WIN32COM_AVAILABLE = True
except ImportError:
    WIN32COM_AVAILABLE = False
    print("Модуль win32com не установлен. PDF не будет создан.")


project_dir = os.path.dirname(os.path.abspath(__file__))
data_dir = os.path.join(project_dir, "data")
os.makedirs(data_dir, exist_ok=True)


def get_week_dates_OATI():
    today = datetime.now()
    monday_last_week = today - timedelta(days=today.weekday() + 7)
    sunday_last_week = monday_last_week + timedelta(days=6)

    monday_last_week_str = monday_last_week.strftime("%d.%m")
    sunday_last_week_str = sunday_last_week.strftime("%d.%m")

    return monday_last_week_str, sunday_last_week_str


def create_ppt_OATI(df, monday_last_week_str, sunday_last_week_str):
    """Создает презентацию ОАТИ"""
    vivod_monday = monday_last_week_str
    vivod_sunday = sunday_last_week_str

    # Список русских месяцев
    months = ['января', 'февраля', 'марта', 'апреля', 'мая', 'июня',
              'июля', 'августа', 'сентября', 'октября', 'ноября', 'декабря']

    # Получаем сегодняшнюю дату
    today = datetime.now()

    # Рассчитываем понедельник и воскресенье прошлой недели
    monday_last_week = today - timedelta(days=today.weekday() + 7)
    sunday_last_week = monday_last_week + timedelta(days=6)

    # Форматируем даты вручную
    monday_last_week_str = f"{monday_last_week.day} {months[monday_last_week.month - 1]} {monday_last_week.year} года"
    sunday_last_week_str = f"{sunday_last_week.day} {months[sunday_last_week.month - 1]} {sunday_last_week.year} года"

    # Формируем текст периода
    period_text = f"период с {monday_last_week_str} по {sunday_last_week_str}"

    # Создаем новую презентацию
    prs = Presentation()

    # Устанавливаем размеры слайда 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Заголовок слайда
    slide_layout = prs.slide_layouts[5]  # Пустой слайд
    slide = prs.slides.add_slide(slide_layout)

    # Добавляем фоновое изображение
    background_image_path = os.path.join(project_dir, '1.png')
    if os.path.exists(background_image_path):
        # Добавляем изображение на задний план
        left = top = Inches(0)
        pic = slide.shapes.add_picture(
            background_image_path, left, top,
            width=prs.slide_width,
            height=prs.slide_height
        )
        # Отправляем изображение на задний план
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(0, pic._element)

    # Заголовок
    title_text = (
        f"Обращения, поступившие в кабинет ОАТИ на портале «Наш город», по которым жители жаловались, "
        f"что их проблема не была устранена и наличие проблемы подтвердилось в результате перепроверки ОАТИ"
    )

    title_box = slide.shapes.add_textbox(Inches(1), Inches(0), Inches(12), Inches(2))  # Заголовок
    title_frame = title_box.text_frame
    title_frame.word_wrap = True  # Включаем перенос слов
    p = title_frame.add_paragraph()
    p.text = title_text
    p.font.size = Pt(18)
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    p_period = title_frame.add_paragraph()  # Новый абзац для периода
    p_period.text = period_text  # Используем период, полученный из get_week_dates
    p_period.font.size = Pt(18)
    p_period.font.name = 'Calibri'
    p_period.alignment = PP_ALIGN.CENTER
    p_period.font.bold = True  # Устанавливаем жирный шрифт

    # Справочная информация
    total_messages = len(df)  # Количество сообщений
    reference_text = f"*Всего опровергнуто за указанный период: {total_messages}"

    reference_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(1))  # Опустили еще ниже
    reference_frame = reference_box.text_frame
    p_ref = reference_frame.add_paragraph()
    p_ref.text = reference_text
    p_ref.font.size = Pt(13)
    p_ref.font.name = 'Calibri'
    p_ref.font.color.rgb = RGBColor(255, 0, 0)  # Красный цвет
    p_ref.alignment = PP_ALIGN.LEFT
    p_ref.font.bold = True  # Устанавливаем жирный шрифт

    # Топ проблемных тем с заданными размерами
    top_issues = df['Проблемная тема'].value_counts().head(3)  # Топ 3 проблемных тем
    top_issues_text = "\n".join([f"{issue} - {count}" for issue, count in top_issues.items()])

    issues_box = slide.shapes.add_textbox(Inches(9.7), Inches(1), Inches(3.5), Inches(1.33))  # Сместили правее
    issues_frame = issues_box.text_frame
    issues_frame.word_wrap = True  # Включаем перенос слов
    p_issues = issues_frame.add_paragraph()
    p_issues.text = "Топ опровергаемых тем:\n" + top_issues_text
    p_issues.font.size = Pt(12)
    p_issues.font.name = 'Calibri'
    p_issues.font.color.rgb = RGBColor(255, 0, 0)  # Красный цвет для текста топ тем
    p_issues.alignment = PP_ALIGN.LEFT
    p_issues.font.bold = True  # Устанавливаем жирный шрифт

    # Получение всех районов
    all_rayon = [
        'АВД ЮВАО', 'Выхино-Жулебино', 'Капотня', 'Кузьминки', 'Лефортово',
        'Люблино', 'Марьино', 'Некрасовка', 'Нижегородский',
        'Печатники', 'Рязанский', 'Текстильщики', 'Южнопортовый'
    ]

    # Считаем количество заявок по районам
    df_count = df['Район'].value_counts().reindex(all_rayon, fill_value=0).reset_index()
    df_count.columns = ['Район', 'Количество заявок']

    # Построение графика
    plt.figure(figsize=(13, 5))  # Делаем график шире и выше
    plt.bar(df_count['Район'], df_count['Количество заявок'], color='#4472c4', width=0.4, zorder=3)

    # Добавляем текст над столбиками
    for index, value in enumerate(df_count['Количество заявок']):
        plt.text(index, value, str(value), ha='center', va='bottom')

    plt.grid(axis='y', linestyle='--', alpha=0.7, zorder=1)
    plt.gca().spines['top'].set_visible(False)
    plt.gca().spines['right'].set_visible(False)
    plt.xticks(rotation=45)
    plt.gca().set_facecolor((1, 1, 1, 0))
    plt.subplots_adjust(bottom=0.24)
    # Сохраняем график с прозрачным фоном
    chart_path = os.path.join(data_dir, 'chart.png')
    plt.savefig(chart_path, transparent=True)
    plt.close()

    # Вставка изображения графика, смещая выше
    left = Inches(-0.7)
    top = Inches(0.5)  # Задаем отступ сверху для графика, больше для повышения его места
    slide.shapes.add_picture(chart_path, left, top, width=Inches(15), height=Inches(6.5))

    # Удаляем временный файл графика
    try:
        os.remove(chart_path)
    except:
        pass

    # Создание сводной таблицы по количеству заявок в каждом районе
    df_pivot = pd.pivot_table(df, values="Номер заявки", index="Район", aggfunc="count")
    df_pivot1 = df_pivot.reset_index()

    # Получение топ-3 районов по количеству заявок
    top_3_rayon = df_pivot1.nlargest(3, 'Номер заявки')
    top_3_rayon_list = top_3_rayon['Район'].tolist()

    # Функция для создания закругленного прямоугольника
    def rectangle1(slide, left, top, widht, height):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, widht, height
        )
        shape.fill.background()  # Полностью прозрачный
        # Устанавливаем цвет обводки
        line = shape.line
        line.color.rgb = RGBColor(255, 0, 0)  # Красный цвет
        line.width = Inches(0.02)  # Ширина обводки
        # Поворачиваем закругленный прямоугольник на 135 градусов
        shape.rotation = 132

    # Добавление закругленных прямоугольников на топ-3 района
    if 'АВД ЮВАО' in top_3_rayon_list:
        rectangle1(slide, Inches(1.28), Inches(5.87), Inches(1.1), Inches(0.3))
    if 'Выхино-Жулебино' in top_3_rayon_list:
        rectangle1(slide, Inches(1.82), Inches(6.15), Inches(1.82), Inches(0.3))
    if 'Капотня' in top_3_rayon_list:
        rectangle1(slide, Inches(2), Inches(6.1), Inches(1.1), Inches(0.27))
    if 'Кузьминки' in top_3_rayon_list:
        rectangle1(slide, Inches(3.9), Inches(5.88), Inches(1.1), Inches(0.27))
    if 'Лефортово' in top_3_rayon_list:
        rectangle1(slide, Inches(4.7), Inches(5.9), Inches(1.1), Inches(0.27))
    if 'Люблино' in top_3_rayon_list:
        rectangle1(slide, Inches(5.67), Inches(5.83), Inches(0.9), Inches(0.27))
    if 'Марьино' in top_3_rayon_list:
        rectangle1(slide, Inches(6.5), Inches(5.83), Inches(0.9), Inches(0.27))
    if 'Некрасовка' in top_3_rayon_list:
        rectangle1(slide, Inches(7.26), Inches(5.92), Inches(1.2), Inches(0.27))
    if 'Нижегородский' in top_3_rayon_list:
        rectangle1(slide, Inches(7.95), Inches(6.1), Inches(1.5), Inches(0.27))
    if 'Печатники' in top_3_rayon_list:
        rectangle1(slide, Inches(8.97), Inches(5.93), Inches(1.1), Inches(0.27))
    if 'Рязанский' in top_3_rayon_list:
        rectangle1(slide, Inches(9.83), Inches(5.91), Inches(1.1), Inches(0.27))
    if 'Текстильщики' in top_3_rayon_list:
        rectangle1(slide, Inches(10.44), Inches(6.04), Inches(1.5), Inches(0.27))
    if 'Южнопортовый' in top_3_rayon_list:
        rectangle1(slide, Inches(11.31), Inches(6.08), Inches(1.55), Inches(0.27))

    # Сохранение презентации
    ppt_filename = f"ОАТИ {vivod_monday}-{vivod_sunday}.pptx"
    ppt_file_path = os.path.join(data_dir, ppt_filename)
    prs.save(ppt_file_path)
    print(f'Презентация сохранена как {ppt_file_path}')

    # Конвертация PPT в PDF
    pdf_filename = f"ОАТИ {vivod_monday}-{vivod_sunday}.pdf"
    pdf_file_path = os.path.join(data_dir, pdf_filename)

    return ppt_file_path


def process_file_OATI(filepath):
    """Обрабатывает файл для создания слайда ОАТИ"""
    try:
        # Загрузите существующий Excel файл
        excel_file = pd.ExcelFile(filepath)
        # Список всех листов
        sheet_names = excel_file.sheet_names
        # Убедитесь, что лист "ОИВ Ответы" существует
        if "ОИВ Ответы" in sheet_names:
            # Считываем только нужный лист
            df = pd.read_excel(filepath, sheet_name="ОИВ Ответы")

            # Фильтрация по столбцу "Округ"
            df = df[df['Округ'] == 'ЮВАО']

            # Фильтрация по столбцу "Ответственный за подготовку ответа"
            df = df[df['Ответственный за подготовку ответа'].isin(['АТИ по ВАО и ЮВАО', 'Дорожная инспекция'])]

            # Фильтрация по столбцу "Категория/Действие ответа"
            df = df[df['Категория/Действие ответа'].isin(['Обещание устранения проблемы', 'Проблема устранена',
                                                          'Проблема устранена до момента проведения проверки'])]

            # Словарь ОИВ 1-го уровня
            mapping_OIV = [
                'ГБУ «Автомобильные дороги ЮВАО»',
                'ГБУ Жилищник Выхино района Выхино-Жулебино города Москвы',
                'ГБУ Жилищник Нижегородского района города Москвы',
                'ГБУ Жилищник района Капотня города Москвы',
                'ГБУ Жилищник района Кузьминки города Москвы',
                'ГБУ Жилищник района Лефортово города Москвы',
                'ГБУ Жилищник района Люблино города Москвы',
                'ГБУ Жилищник района Марьино города Москвы',
                'ГБУ Жилищник района Некрасовка города Москвы',
                'ГБУ Жилищник района Печатники города Москвы',
                'ГБУ Жилищник района Текстильщики города Москвы',
                'ГБУ Жилищник Рязанского района города Москвы',
                'ГБУ Жилищник Южнопортового района города Москвы',
                'Управа Выхино-Жулебино',
                'Управа Нижегородский',
                'Управа Капотня',
                'Управа Кузьминки',
                'Управа Лефортово',
                'Управа Люблино',
                'Управа Марьино',
                'Управа Некрасовка',
                'Управа Печатники',
                'Управа Текстильщики',
                'Управа Рязанский',
                'Управа Южнопортовый'
            ]

            # Фильтрация по столбцу "Ответственный ОИВ первого уровня"
            df = df[df['Ответственный ОИВ первого уровня'].isin(mapping_OIV)]

            # Словарь ОИВ 1-го уровня по отношению к районю
            responsible_mapping = {
                'ГБУ «Автомобильные дороги ЮВАО»': 'АВД ЮВАО',
                'ГБУ Жилищник Выхино района Выхино-Жулебино города Москвы': 'Выхино-Жулебино',
                'Управа Выхино-Жулебино': 'Выхино-Жулебино',
                'ГБУ Жилищник Нижегородского района города Москвы': 'Нижегородский',
                'Управа Нижегородский': 'Нижегородский',
                'ГБУ Жилищник района Капотня города Москвы': 'Капотня',
                'Управа Капотня': 'Капотня',
                'ГБУ Жилищник района Кузьминки города Москвы': 'Кузьминки',
                'Управа Кузьминки': 'Кузьминки',
                'ГБУ Жилищник района Лефортово города Москвы': 'Лефортово',
                'Управа Лефортово': 'Лефортово',
                'ГБУ Жилищник района Люблино города Москвы': 'Люблино',
                'Управа Люблино': 'Люблино',
                'ГБУ Жилищник района Марьино города Москвы': 'Марьино',
                'Управа Марьино': 'Марьино',
                'ГБУ Жилищник района Некрасовка города Москвы': 'Некрасовка',
                'Управа Некрасовка': 'Некрасовка',
                'ГБУ Жилищник района Печатники города Москвы': 'Печатники',
                'Управа Печатники': 'Печатники',
                'ГБУ Жилищник района Текстильщики города Москвы': 'Текстильщики',
                'Управа Текстильщики': 'Текстильщики',
                'ГБУ Жилищник Рязанского района города Москвы': 'Рязанский',
                'Управа Рязанский': 'Рязанский',
                'ГБУ Жилищник Южнопортового района города Москвы': 'Южнопортовый',
                'Управа Южнопортовый': 'Южнопортовый'
            }

            # Замена района по ОИВ 1-го уровня
            df['Район'] = df['Ответственный ОИВ первого уровня'].map(responsible_mapping)

            # Получаем даты для названия файла
            monday_last_week_str, sunday_last_week_str = get_week_dates_OATI()

            # Создание презентации
            ppt_path = create_ppt_OATI(df, monday_last_week_str, sunday_last_week_str)

            # Подготовка статистики для сообщения
            # Всего сообщений
            total_messages = len(df)

            # Устранено
            ustraneno = len(df[df['Категория/Действие ответа'] == 'Проблема устранена'])

            # Переносов
            perenosov = len(df[df['Категория/Действие ответа'] == 'Обещание устранения проблемы'])

            # Самый долгий перенос
            # Фильтруем только строки с переносами
            perenos_df = df[df['Категория/Действие ответа'] == 'Обещание устранения проблемы'].copy()

            if not perenos_df.empty:
                # Функция для поиска даты в тексте ответа
                def find_date_in_text(text):
                    if pd.isna(text):
                        return None

                    # Пробуем найти дату в формате DD.MM.YYYY
                    date_pattern = r'\b(\d{1,2}\.\d{1,2}\.\d{4})\b'
                    matches = re.findall(date_pattern, str(text))

                    if matches:
                        # Берем первую найденную дату
                        try:
                            date_str = matches[0]
                            # Конвертируем в datetime
                            return datetime.strptime(date_str, '%d.%m.%Y')
                        except:
                            return None

                    return None

                # Применяем функцию к столбцу "Текст ответа"
                perenos_df['found_date'] = perenos_df['Текст ответа'].apply(find_date_in_text)

                # Удаляем строки, где дата не найдена
                perenos_df = perenos_df.dropna(subset=['found_date'])

                if not perenos_df.empty:
                    # Находим самую позднюю дату (самый долгий перенос)
                    latest_date_row = perenos_df.loc[perenos_df['found_date'].idxmax()]
                    longest_perenos_date = latest_date_row['found_date'].strftime('%d.%m.%Y')
                    longest_perenos_id = latest_date_row['Номер заявки']
                else:
                    longest_perenos_date = "не найдено"
                    longest_perenos_id = "не найден"
            else:
                longest_perenos_date = "нет переносов"
                longest_perenos_id = "нет переносов"

            # Формируем сообщение
            message = f"Всего - {total_messages}\n"
            message += f"Устранено - {ustraneno}\n"
            message += f"Переносов - {perenosov}\n"
            message += f"Самый долгий перенос - {longest_perenos_date} ({longest_perenos_id})"

            excel_file.close()

            return ppt_path, message
        else:
            raise Exception('Лист "ОИВ Ответы" не найден в файле.')

    except Exception as e:
        raise Exception(f'Ошибка обработки файла ОАТИ: {str(e)}')
